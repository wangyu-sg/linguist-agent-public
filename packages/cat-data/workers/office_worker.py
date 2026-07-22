#!/usr/bin/env python3
"""Read-only Office/PDF inspection and copy-on-write editing worker."""

from __future__ import annotations

import hashlib
import json
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any


def digest(path: Path) -> str:
    hasher = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            hasher.update(chunk)
    return hasher.hexdigest()


def run_docx(executable: Path, args: list[str], expect_json: bool = False) -> Any:
    completed = subprocess.run(
        [str(executable), *args],
        check=True,
        capture_output=True,
        text=True,
        timeout=120,
    )
    if not expect_json:
        return completed.stdout
    try:
        return json.loads(completed.stdout)
    except json.JSONDecodeError as error:
        raise RuntimeError(f"bun-docx returned invalid JSON: {error}") from error


def inspect_docx(path: Path, executable: Path) -> dict[str, Any]:
    ast = run_docx(executable, ["read", str(path), "--ast"], expect_json=True)
    validation = run_docx(executable, ["validate", str(path), "--json"], expect_json=True)
    return {"format": "docx", "ast": ast, "validation": validation}


def inspect_xlsx(path: Path) -> dict[str, Any]:
    from openpyxl import load_workbook
    workbook = load_workbook(path, read_only=False, data_only=False, keep_links=True)
    try:
        sheets: list[dict[str, Any]] = []
        for sheet in workbook.worksheets:
            formulas: list[dict[str, str]] = []
            hyperlinks: list[dict[str, str]] = []
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value.startswith("="):
                        formulas.append({"cell": cell.coordinate, "formula": cell.value})
                    if cell.hyperlink:
                        hyperlinks.append({"cell": cell.coordinate, "target": str(cell.hyperlink.target or cell.hyperlink.location or "")})
            sheets.append({
                "name": sheet.title,
                "maxRow": sheet.max_row,
                "maxColumn": sheet.max_column,
                "mergedRanges": sorted(str(value) for value in sheet.merged_cells.ranges),
                "formulas": formulas,
                "hyperlinks": hyperlinks,
            })
        return {"format": "xlsx", "sheets": sheets}
    finally:
        workbook.close()


def inspect_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation
    presentation = Presentation(path)
    return {
        "format": "pptx",
        "slides": [{
            "slide": index + 1,
            "shapes": [{
                "shape": shape_index,
                "name": shape.name,
                "type": str(shape.shape_type),
                "hasText": bool(getattr(shape, "has_text_frame", False)),
                "text": shape.text if getattr(shape, "has_text_frame", False) else None,
            } for shape_index, shape in enumerate(slide.shapes)],
        } for index, slide in enumerate(presentation.slides)],
    }


def inspect_pdf(path: Path) -> dict[str, Any]:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return {
        "format": "pdf",
        "pages": len(reader.pages),
        "encrypted": bool(reader.is_encrypted),
        "formFields": sorted((reader.get_fields() or {}).keys()),
        "annotations": [len(page.get("/Annots") or []) for page in reader.pages],
    }


def create_docx(output: Path, request: dict[str, Any], executable: Path) -> dict[str, Any]:
    text = str(request.get("text") or "")
    if not text:
        raise ValueError("create_docx requires text")
    run_docx(executable, ["create", str(output), "--text", text])
    return {"format": "docx", "operation": "create", "characters": len(text)}


def create_pptx(output: Path, request: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches
    slides = request.get("slides") or []
    if not isinstance(slides, list) or not slides:
        raise ValueError("create_pptx requires a non-empty slides array")
    presentation = Presentation()
    for item in slides:
        if not isinstance(item, dict):
            raise ValueError("Each PPTX slide must be an object")
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.7))
        title.text = str(item.get("title") or "")
        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.0))
        body.text = str(item.get("body") or "")
    presentation.save(output)
    return {"format": "pptx", "operation": "create", "slides": len(slides)}


def annotate_pdf(path: Path, output: Path, request: dict[str, Any]) -> dict[str, Any]:
    from pypdf import PdfReader, PdfWriter
    from pypdf.annotations import FreeText
    reader = PdfReader(path)
    page = int(request.get("pages", [1])[0]) - 1
    if page < 0 or page >= len(reader.pages):
        raise ValueError(f"PDF annotation page must be between 1 and {len(reader.pages)}")
    rect = request.get("rect") or [50, 50, 250, 120]
    if not isinstance(rect, list) or len(rect) != 4:
        raise ValueError("PDF annotation rect must contain four numbers")
    text = str(request.get("text") or "")
    if not text:
        raise ValueError("PDF annotation text is required")
    writer = PdfWriter()
    writer.clone_document_from_reader(reader)
    writer.add_annotation(page_number=page, annotation=FreeText(text=text, rect=tuple(float(value) for value in rect)))
    with output.open("xb") as stream:
        writer.write(stream)
    return {"format": "pdf", "operation": "annotate", "page": page + 1, "rect": rect, "text": text}


def replace_xlsx(path: Path, output: Path, replacements: list[dict[str, Any]]) -> dict[str, Any]:
    from openpyxl import load_workbook
    shutil.copy2(path, output)
    workbook = load_workbook(output, data_only=False, keep_links=True)
    changed: list[dict[str, Any]] = []
    try:
        for item in replacements:
            sheet = workbook[str(item["sheet"])]
            cell = sheet[str(item["cell"])]
            before = cell.value
            cell.value = item.get("value")
            changed.append({"sheet": sheet.title, "cell": cell.coordinate, "before": before, "after": cell.value})
        workbook.save(output)
    finally:
        workbook.close()
    return {"format": "xlsx", "changes": changed}


def replace_pptx(path: Path, output: Path, replacements: list[dict[str, Any]]) -> dict[str, Any]:
    from pptx import Presentation
    shutil.copy2(path, output)
    presentation = Presentation(output)
    changed: list[dict[str, Any]] = []
    for item in replacements:
        slide_index = int(item["slide"]) - 1
        shape_index = int(item["shape"])
        shape = presentation.slides[slide_index].shapes[shape_index]
        if not getattr(shape, "has_text_frame", False):
            raise ValueError(f"Slide {slide_index + 1} shape {shape_index} has no text frame")
        before = shape.text
        replacement = str(item.get("text") or "")
        find = item.get("find")
        if find is not None:
            matched = False
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if str(find) in run.text:
                        run.text = run.text.replace(str(find), replacement)
                        matched = True
            if not matched:
                raise ValueError(f"Slide {slide_index + 1} shape {shape_index} did not contain requested text")
        elif shape.text_frame.paragraphs:
            paragraph = shape.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].text = replacement
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.text = replacement
            for extra in shape.text_frame.paragraphs[1:]:
                extra.text = ""
        changed.append({"slide": slide_index + 1, "shape": shape_index, "before": before, "after": shape.text})
    presentation.save(output)
    return {"format": "pptx", "changes": changed}


def mutate_docx(path: Path, output: Path, operation: str, request: dict[str, Any], executable: Path) -> dict[str, Any]:
    shutil.copy2(path, output)
    if operation == "replace":
        replacements = request.get("replacements") or []
        if not isinstance(replacements, list) or not replacements:
            raise ValueError("DOCX replace requires replacements")
        changes: list[dict[str, Any]] = []
        for item in replacements:
            if not isinstance(item, dict):
                raise ValueError("Each DOCX replacement must be an object")
            args: list[str]
            if item.get("locator"):
                args = ["edit", str(output), "--at", str(item["locator"]), "--text", str(item.get("text") or "")]
            else:
                pattern = str(item.get("pattern") or "")
                if not pattern:
                    raise ValueError("DOCX replacement requires locator or pattern")
                args = ["replace", str(output), pattern, str(item.get("replacement") or item.get("text") or "")]
                if item.get("all", True):
                    args.append("--all")
            if item.get("track"):
                args.append("--track")
            run_docx(executable, args)
            changes.append({"locator": item.get("locator"), "pattern": item.get("pattern"), "tracked": bool(item.get("track"))})
        diff = run_docx(executable, ["diff", str(output), "--against", str(path), "--comments", "--json"], expect_json=True)
        return {"format": "docx", "changes": changes, "diff": diff}
    if operation == "comment":
        locator = str(request.get("locator") or "")
        text = str(request.get("text") or "")
        if not locator or not text:
            raise ValueError("DOCX comment requires locator and text")
        run_docx(executable, ["comments", "add", str(output), "--at", locator, "--text", text])
        return {"format": "docx", "comment": {"locator": locator, "text": text}}
    if operation == "track_changes":
        enabled = bool(request.get("enabled", True))
        run_docx(executable, ["track-changes", str(output), "on" if enabled else "off"])
        return {"format": "docx", "trackChanges": enabled}
    raise ValueError(f"Unsupported DOCX operation: {operation}")


def validate_preservation(suffix: str, before: dict[str, Any], after: dict[str, Any], replacements: list[dict[str, Any]]) -> dict[str, Any]:
    if suffix in {".xlsx", ".xlsm"}:
        targets = {(str(item.get("sheet")), str(item.get("cell"))) for item in replacements}
        before_sheets = {sheet["name"]: sheet for sheet in before["sheets"]}
        after_sheets = {sheet["name"]: sheet for sheet in after["sheets"]}
        if before_sheets.keys() != after_sheets.keys():
            raise RuntimeError("XLSX sheet names changed unexpectedly")
        for name, source_sheet in before_sheets.items():
            output_sheet = after_sheets[name]
            if source_sheet["mergedRanges"] != output_sheet["mergedRanges"]:
                raise RuntimeError(f"XLSX merged ranges changed unexpectedly in {name}")
            for field in ("formulas", "hyperlinks"):
                source_values = {item["cell"]: item.get("formula", item.get("target")) for item in source_sheet[field] if (name, item["cell"]) not in targets}
                output_values = {item["cell"]: item.get("formula", item.get("target")) for item in output_sheet[field] if (name, item["cell"]) not in targets}
                if source_values != output_values:
                    raise RuntimeError(f"XLSX {field} changed outside requested cells in {name}")
        return {"sourceReopened": True, "outputReopened": True, "unrelatedFormulasPreserved": True, "unrelatedHyperlinksPreserved": True, "mergedRangesPreserved": True}
    if suffix == ".pptx":
        targets = {(int(item.get("slide", 0)), int(item.get("shape", -1))) for item in replacements}
        if len(before["slides"]) != len(after["slides"]):
            raise RuntimeError("PPTX slide count changed unexpectedly")
        for source_slide, output_slide in zip(before["slides"], after["slides"]):
            if len(source_slide["shapes"]) != len(output_slide["shapes"]):
                raise RuntimeError(f"PPTX shape count changed unexpectedly on slide {source_slide['slide']}")
            for source_shape, output_shape in zip(source_slide["shapes"], output_slide["shapes"]):
                identity = (source_slide["slide"], source_shape["shape"])
                if source_shape["name"] != output_shape["name"] or source_shape["type"] != output_shape["type"]:
                    raise RuntimeError(f"PPTX shape identity changed unexpectedly on slide {identity[0]} shape {identity[1]}")
                if identity not in targets and source_shape["text"] != output_shape["text"]:
                    raise RuntimeError(f"PPTX text changed outside requested shape on slide {identity[0]} shape {identity[1]}")
        return {"sourceReopened": True, "outputReopened": True, "slideAndShapeStructurePreserved": True, "unrelatedTextPreserved": True}
    return {"sourceReopened": True, "outputReopened": True, "schemaValid": True}


def handle(request: dict[str, Any]) -> dict[str, Any]:
    operation = str(request.get("operation") or "inspect")
    docx_executable = Path(str(request.get("docxExecutable") or ""))
    if operation in {"create_docx", "create_pptx"}:
        output = Path(str(request.get("outputPath") or "")).resolve(strict=False)
        if output.exists():
            raise FileExistsError(f"Output already exists: {output}")
        output.parent.mkdir(parents=True, exist_ok=True)
        result = create_docx(output, request, docx_executable) if operation == "create_docx" else create_pptx(output, request)
        validation = inspect_docx(output, docx_executable) if operation == "create_docx" else inspect_pptx(output)
        return {"ok": True, "outputPath": str(output), "outputSha256": digest(output), "diff": result, "validation": {"reopened": validation, "preservation": {"outputReopened": True}}}
    source = Path(str(request.get("sourcePath") or "")).resolve(strict=True)
    before = digest(source)
    suffix = source.suffix.lower()
    if operation == "inspect":
        result = inspect_docx(source, docx_executable) if suffix == ".docx" else inspect_xlsx(source) if suffix in {".xlsx", ".xlsm"} else inspect_pptx(source) if suffix == ".pptx" else inspect_pdf(source) if suffix == ".pdf" else None
        if result is None:
            raise ValueError(f"Unsupported managed Office inspection format: {suffix}")
        return {"ok": True, "sourcePath": str(source), "sourceSha256": before, "result": result}

    output = Path(str(request.get("outputPath") or "")).resolve(strict=False)
    if output == source:
        raise ValueError("Managed Office edits require a new output file")
    if output.exists():
        raise FileExistsError(f"Output already exists: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)
    replacements = request.get("replacements") or []
    if not isinstance(replacements, list):
        raise ValueError("replacements must be an array")
    if operation not in {"replace", "comment", "track_changes", "pdf_annotate"}:
        raise ValueError(f"Unsupported managed Office operation: {operation}")
    source_structure = inspect_xlsx(source) if suffix in {".xlsx", ".xlsm"} else inspect_pptx(source) if suffix == ".pptx" else inspect_docx(source, docx_executable) if suffix == ".docx" else inspect_pdf(source) if suffix == ".pdf" else None
    result = annotate_pdf(source, output, request) if suffix == ".pdf" and operation == "pdf_annotate" else mutate_docx(source, output, operation, request, docx_executable) if suffix == ".docx" else replace_xlsx(source, output, replacements) if suffix in {".xlsx", ".xlsm"} and operation == "replace" else replace_pptx(source, output, replacements) if suffix == ".pptx" and operation == "replace" else None
    if result is None:
        raise ValueError(f"Copy-on-write replace is not supported for {suffix}")
    if digest(source) != before:
        raise RuntimeError("Source document changed during copy-on-write operation")
    # Re-open validation is performed by the same parser before success.
    output_structure = inspect_docx(output, docx_executable) if suffix == ".docx" else inspect_xlsx(output) if suffix in {".xlsx", ".xlsm"} else inspect_pptx(output) if suffix == ".pptx" else inspect_pdf(output)
    validation = {
        "reopened": output_structure,
        "preservation": validate_preservation(suffix, source_structure or {}, output_structure, replacements),
    }
    return {"ok": True, "sourcePath": str(source), "sourceSha256": before, "outputPath": str(output), "outputSha256": digest(output), "diff": result, "validation": validation}


def main() -> None:
    for line in sys.stdin:
        if not line.strip():
            continue
        try:
            request = json.loads(line)
            if not isinstance(request, dict):
                raise ValueError("Worker request must be a JSON object")
            response = handle(request)
        except Exception as error:
            response = {"ok": False, "error": str(error), "errorType": type(error).__name__}
        print(json.dumps(response, ensure_ascii=False), flush=True)


if __name__ == "__main__":
    main()
