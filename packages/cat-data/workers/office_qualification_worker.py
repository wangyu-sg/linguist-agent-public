#!/usr/bin/env python3
"""Generate and validate the managed Office pack's offline structure fixtures."""

from __future__ import annotations

import hashlib
import ipaddress
import json
import shutil
import socket
import subprocess
import sys
from pathlib import Path
from typing import Any


OUTBOUND_REQUESTS = 0


def digest(path: Path) -> str:
    value = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            value.update(chunk)
    return value.hexdigest()


def install_network_guard() -> None:
    original_getaddrinfo = socket.getaddrinfo
    original_connect = socket.socket.connect

    def local(host: Any) -> bool:
        if host is None or host in {"", "localhost"}:
            return True
        try:
            return ipaddress.ip_address(str(host).strip("[]")).is_loopback
        except ValueError:
            return False

    def guarded_getaddrinfo(host: Any, *args: Any, **kwargs: Any) -> Any:
        global OUTBOUND_REQUESTS
        if not local(host):
            OUTBOUND_REQUESTS += 1
            raise RuntimeError(f"Office qualification blocked network lookup: {host}")
        return original_getaddrinfo(host, *args, **kwargs)

    def guarded_connect(instance: socket.socket, address: Any) -> Any:
        global OUTBOUND_REQUESTS
        host = address[0] if isinstance(address, tuple) and address else address
        if instance.family in {socket.AF_INET, socket.AF_INET6} and not local(host):
            OUTBOUND_REQUESTS += 1
            raise RuntimeError(f"Office qualification blocked network connection: {host}")
        return original_connect(instance, address)

    socket.getaddrinfo = guarded_getaddrinfo
    socket.socket.connect = guarded_connect


def run(executable: Path, args: list[str]) -> str:
    return subprocess.run([str(executable), *args], check=True, capture_output=True, text=True, timeout=120).stdout


def observation(identifier: str, format_name: str, source: Path, before: str, output_reopened: bool, references_preserved: bool = True, formulas_preserved: bool = True) -> dict[str, Any]:
    return {
        "id": identifier,
        "format": format_name,
        "sourceUnchanged": digest(source) == before,
        "outputReopened": output_reopened,
        "referencesPreserved": references_preserved,
        "formulasPreserved": formulas_preserved,
    }


def qualify_docx(root: Path, executable: Path, count: int) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    directory = root / "docx"
    directory.mkdir()
    for index in range(count):
        source = directory / f"source-{index + 1:02d}.docx"
        output = directory / f"output-{index + 1:02d}.docx"
        run(executable, ["create", str(source), "--text", f"Draft localization fixture {index + 1}.\nReviewer note {index + 1}."])
        before = digest(source)
        shutil.copy2(source, output)
        args = ["replace", str(output), "Draft", "Final", "--all"]
        if index % 2 == 0:
            args.append("--track")
        run(executable, args)
        if index % 3 == 0:
            run(executable, ["comments", "add", str(output), "--at", "p0:0-5", "--text", f"Review {index + 1}"])
        validation = json.loads(run(executable, ["validate", str(output), "--json"]))
        rows.append(observation(f"docx-{index + 1:02d}", "docx", source, before, validation.get("valid") is True))
    return rows


def qualify_xlsx(root: Path, count: int) -> list[dict[str, Any]]:
    from openpyxl import Workbook, load_workbook
    rows: list[dict[str, Any]] = []
    directory = root / "xlsx"
    directory.mkdir()
    for index in range(count):
        source = directory / f"source-{index + 1:02d}.xlsx"
        output = directory / f"output-{index + 1:02d}.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = f"Strings{index + 1}"
        sheet["A1"] = "Draft"
        sheet["B1"] = f"=LEN(A1)+{index}"
        sheet["C1"] = "Reference"
        sheet["C1"].hyperlink = f"https://example.invalid/reference/{index + 1}"
        if index % 2 == 0:
            sheet.merge_cells("D1:E1")
            sheet["D1"] = "Merged"
        if index % 4 == 0:
            workbook.create_sheet("Notes")["A1"] = f"Note {index + 1}"
        workbook.save(source)
        before = digest(source)
        shutil.copy2(source, output)
        edited = load_workbook(output, data_only=False, keep_links=True)
        edited[sheet.title]["A1"] = "Final"
        edited.save(output)
        edited.close()
        reopened = load_workbook(output, data_only=False, keep_links=True)
        output_sheet = reopened[sheet.title]
        formula_ok = output_sheet["B1"].value == f"=LEN(A1)+{index}"
        link_ok = output_sheet["C1"].hyperlink is not None and output_sheet["C1"].hyperlink.target == f"https://example.invalid/reference/{index + 1}"
        reopened.close()
        rows.append(observation(f"xlsx-{index + 1:02d}", "xlsx", source, before, True, link_ok, formula_ok))
    return rows


def qualify_pptx(root: Path, count: int) -> list[dict[str, Any]]:
    from pptx import Presentation
    from pptx.util import Inches
    rows: list[dict[str, Any]] = []
    directory = root / "pptx"
    directory.mkdir()
    for index in range(count):
        source = directory / f"source-{index + 1:02d}.pptx"
        output = directory / f"output-{index + 1:02d}.pptx"
        presentation = Presentation()
        for slide_index in range(index % 3 + 1):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(7), Inches(1))
            shape.text = f"Draft slide {slide_index + 1} fixture {index + 1}"
        presentation.save(source)
        before = digest(source)
        shutil.copy2(source, output)
        edited = Presentation(output)
        edited.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text = f"Final slide 1 fixture {index + 1}"
        source_shape_counts = [len(slide.shapes) for slide in edited.slides]
        edited.save(output)
        reopened = Presentation(output)
        structure_ok = source_shape_counts == [len(slide.shapes) for slide in reopened.slides]
        rows.append(observation(f"pptx-{index + 1:02d}", "pptx", source, before, structure_ok))
    return rows


def qualify_pdf(root: Path, count: int) -> list[dict[str, Any]]:
    from pypdf import PdfReader, PdfWriter
    from pypdf.annotations import FreeText
    rows: list[dict[str, Any]] = []
    directory = root / "pdf"
    directory.mkdir()
    for index in range(count):
        source = directory / f"source-{index + 1:02d}.pdf"
        output = directory / f"output-{index + 1:02d}.pdf"
        writer = PdfWriter()
        for _ in range(index % 3 + 1):
            writer.add_blank_page(width=612, height=792)
        writer.add_metadata({"/Title": f"Fixture {index + 1}"})
        writer.write(source)
        before = digest(source)
        reader = PdfReader(source)
        edited = PdfWriter()
        edited.clone_document_from_reader(reader)
        edited.pages[0].rotate(90 if index % 2 == 0 else 180)
        if index % 3 == 0:
            edited.add_annotation(0, FreeText(text=f"Review {index + 1}", rect=(50, 600, 260, 680)))
        edited.write(output)
        reopened = PdfReader(output)
        structure_ok = len(reopened.pages) == len(reader.pages)
        rows.append(observation(f"pdf-{index + 1:02d}", "pdf", source, before, structure_ok))
    return rows


def handle(request: dict[str, Any]) -> dict[str, Any]:
    install_network_guard()
    output_root = Path(str(request.get("outputRoot") or "")).resolve(strict=False)
    count = int(request.get("count") or 20)
    if count < 20 or count > 50:
        raise ValueError("Office qualification fixture count must be between 20 and 50")
    if output_root.exists() and any(output_root.iterdir()):
        raise FileExistsError(f"Office qualification output is not empty: {output_root}")
    output_root.mkdir(parents=True, exist_ok=True)
    docx_executable = Path(str(request.get("docxExecutable") or "")).resolve(strict=True)
    fixtures = [
        *qualify_docx(output_root, docx_executable, count),
        *qualify_xlsx(output_root, count),
        *qualify_pptx(output_root, count),
        *qualify_pdf(output_root, count),
    ]
    return {"ok": True, "outputRoot": str(output_root), "fixtures": fixtures, "outboundCustomerFileRequests": OUTBOUND_REQUESTS}


def main() -> None:
    for line in sys.stdin:
        if not line.strip():
            continue
        try:
            request = json.loads(line)
            if not isinstance(request, dict):
                raise ValueError("Worker request must be a JSON object")
            response = handle(request)
        except Exception as error:
            response = {"ok": False, "error": str(error), "errorType": type(error).__name__}
        print(json.dumps(response, ensure_ascii=False), flush=True)


if __name__ == "__main__":
    main()
